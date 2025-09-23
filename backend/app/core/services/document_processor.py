"""
文書処理モジュール
PDF読み込み、テキスト抽出、チャンク分割などの
文書処理に関する機能を提供するモジュール
"""

import io

import pypdf
from langchain.text_splitter import RecursiveCharacterTextSplitter

from ..config import settings


class DocumentProcessor:
    """文書処理クラス
    PDF文書の読み込み、テキスト抽出、チャンク分割を行う
    """

    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.max_chunk_size,
            chunk_overlap=settings.chunk_overlap,
            separators=["\\n\\n", "\\n", " ", ""],
            length_function=len,
        )

    def split_text(
        self,
        text: str,
        chunk_size: int | None = None,
        chunk_overlap: int | None = None,
    ) -> list[str]:
        """テキストをチャンクに分割"""
        if chunk_size or chunk_overlap:
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size or settings.max_chunk_size,
                chunk_overlap=chunk_overlap or settings.chunk_overlap,
                separators=["\\n\\n", "\\n", " ", ""],
                length_function=len,
            )
            return splitter.split_text(text)
        return self.text_splitter.split_text(text)

    def extract_text_from_pdf(self, data: bytes) -> str:
        """PDFからテキスト抽出"""
        try:
            reader = pypdf.PdfReader(io.BytesIO(data))
            if len(reader.pages) == 0:
                raise ValueError("PDFにページが含まれていません")
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            if not text_parts:
                raise ValueError("PDFからテキストを抽出できませんでした")
            return "\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"PDF読み込みエラー: {str(e)}")

    def extract_text_from_txt_bytes(self, data: bytes) -> str:
        """TXT(プレーンテキスト)からテキスト抽出"""
        try:
            if data.startswith(b"\xef\xbb\xbf"):
                enc = "utf-8-sig"
            else:
                enc = "utf-8"
            return data.decode(enc, errors="replace")
        except Exception as e:
            raise ValueError(f"TXT読み込みエラー: {str(e)}")

    def extract_text_from_docx_bytes(self, data: bytes) -> str:
        """DOCX(Word)からテキスト抽出"""
        try:
            from docx import Document

            buf = io.BytesIO(data)
            doc = Document(buf)
            parts: list[str] = []
            for p in doc.paragraphs:
                if p.text:
                    parts.append(p.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            parts.append(t)
            return "\n".join(parts)
        except Exception as e:
            raise ValueError(f"DOCX読み込みエラー: {str(e)}")

    def extract_text_from_pptx_bytes(self, data: bytes) -> str:
        """PPTX(Power Point)からテキスト抽出"""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = (
                                "".join(run.text for run in paragraph.runs)
                                or paragraph.text
                            )
                            t = text.strip()
                            if t:
                                parts.append(t)
                if getattr(slide, "has_notes_slide", False) and slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame
                    if notes:
                        for paragraph in notes.paragraphs:
                            text = (
                                "".join(run.text for run in paragraph.runs)
                                or paragraph.text
                            )
                            t = text.strip()
                            if t:
                                parts.append(t)
            return "\n".join(parts)
        except Exception as e:
            raise ValueError(f"PPTX読み込みエラー: {str(e)}")

    def extract_text_from_xlsx_bytes(self, data: bytes) -> str:
        """XLSX(Excel)からテキスト抽出"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(
                io.BytesIO(data), data_only=True, read_only=True
            )
            parts: list[str] = []
            for ws in wb.worksheets:
                parts.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [
                        str(c).strip() for c in row if c is not None and str(c).strip()
                    ]
                    if cells:
                        parts.append("\t".join(cells))
            return "\n".join(parts)
        except Exception as e:
            raise ValueError(f"XLSX読み込みエラー: {str(e)}")
